import os
import matplotlib.pyplot as plt
import ezdxf
from ezdxf.addons.drawing import RenderContext, Frontend
from ezdxf.addons.drawing.matplotlib import MatplotlibBackend
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM
from pptx import Presentation
from pptx.util import Inches

def convert_dxf_to_svg(dxf_path, svg_path):
    """
    Converts a DXF file to an SVG file using Matplotlib and ezdxf.
    
    :param dxf_path: Path to the input DXF file.
    :param svg_path: Path to save the output SVG file.
    """
    try:
        # Ensure output folder exists
        os.makedirs(os.path.dirname(svg_path), exist_ok=True)

        # Load the DXF file
        doc = ezdxf.readfile(dxf_path)

        # Create a figure and backend for Matplotlib
        fig = plt.figure()
        out = MatplotlibBackend(fig.add_axes([0, 0, 1, 1]))

        # Render the DXF to the Matplotlib backend
        Frontend(RenderContext(doc), out).draw_layout(doc.modelspace(), finalize=True)

        # Save the result as an SVG
        fig.savefig(svg_path, format="svg")
        print(f"DXF successfully converted to SVG and saved at {svg_path}")
    except Exception as e:
        print(f"Error during DXF to SVG conversion: {e}")
        # If conversion fails, create an empty SVG file as a fallback
        with open(svg_path, "w") as f:
            f.write("<svg xmlns='http://www.w3.org/2000/svg' width='100' height='100'></svg>")
        print(f"Fallback: Created empty SVG file at {svg_path}")

def convert_svg_to_emf(svg_path, emf_path):
    """
    Converts an SVG file to EMF format using svglib and reportlab.
    
    :param svg_path: Path to the input SVG file.
    :param emf_path: Path to save the output EMF file.
    """
    try:
        # Ensure file exists before attempting conversion
        if not os.path.exists(svg_path):
            print(f"SVG file {svg_path} does not exist.")
            return
        
        # Load SVG file
        drawing = svg2rlg(svg_path)

        # Convert to EMF and save
        renderPM.drawToFile(drawing, emf_path, fmt="EMF")
        print(f"SVG successfully converted to EMF and saved at {emf_path}")
    except Exception as e:
        print(f"Error during SVG to EMF conversion: {e}")

def add_emf_to_pptx(emf_path, pptx_path):
    """
    Adds an EMF file to a PowerPoint presentation.
    
    :param emf_path: Path to the input EMF file.
    :param pptx_path: Path to save the output PowerPoint file.
    """
    try:
        # Create a PowerPoint presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # Blank slide
        slide = prs.slides.add_slide(slide_layout)

        # Add EMF image to the slide
        slide.shapes.add_picture(emf_path, Inches(1), Inches(1), Inches(5), Inches(5))

        # Save the presentation
        prs.save(pptx_path)
        print(f"EMF successfully inserted into PowerPoint and saved at {pptx_path}")
    except Exception as e:
        print(f"Error during PowerPoint creation: {e}")

# File paths
dxf_file = "./test/sample_dxf.dxf"  # Replace with your DXF file path
svg_file = "./test/sample_dxf.svg"  # Replace with your desired SVG output path
emf_file = "./test/sample_dxf.emf"  # Replace with your desired EMF output path
pptx_file = "./test/sample_dxf.pptx"  # Replace with your desired PowerPoint output path

# Convert DXF to SVG
convert_dxf_to_svg(dxf_file, svg_file)

# Convert SVG to EMF
convert_svg_to_emf(svg_file, emf_file)

# Add EMF to PowerPoint
add_emf_to_pptx(emf_file, pptx_file)
